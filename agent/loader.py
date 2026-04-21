"""loader.py - Multi-format Document Loader v3.0"""
import os, io, time, logging, re
from pathlib import Path
from typing import List, Dict, Any, Optional, Tuple
from dataclasses import dataclass
from threading import Lock
from concurrent.futures import ThreadPoolExecutor

logger = logging.getLogger(__name__)

@dataclass
class Document:
    page_content: str
    metadata: Dict[str, Any]
    chunk_id: Optional[str] = None

    def to_dict(self) -> Dict[str, Any]:
        return {
            'page_content': self.page_content,
            'metadata': self.metadata,
            'chunk_id': self.chunk_id
        }

@dataclass
class LoadResult:
    success: bool
    documents: List[Document]
    error: Optional[str] = None
    file_type: str = ''
    page_count: int = 0
    char_count: int = 0
    processing_time: float = 0.0

    def to_dict(self) -> Dict[str, Any]:
        return {
            'success': self.success,
            'document_count': len(self.documents),
            'error': self.error,
            'file_type': self.file_type,
            'page_count': self.page_count,
            'char_count': self.char_count,
            'processing_time': round(self.processing_time, 3)
        }

class DocumentLoader:
    _instance = None
    _lock = Lock()

    def __new__(cls):
        if cls._instance is None:
            with cls._lock:
                if cls._instance is None:
                    cls._instance = super().__new__(cls)
                    cls._instance._init()
        return cls._instance

    def _init(self) -> None:
        self._supported_formats = {
            '.pdf': self.load_pdf,
            '.txt': self.load_txt,
            '.md': self.load_md,
            '.docx': self.load_docx,
            '.xlsx': self.load_xlsx,
            '.pptx': self.load_pptx,
            '.csv': self.load_csv,
            '.html': self.load_html,
            '.htm': self.load_html,
            '.json': self.load_json,
            '.xml': self.load_xml,
            '.eml': self.load_eml,
            '.msg': self.load_msg,
            '.png': self.load_image,
            '.jpg': self.load_image,
            '.jpeg': self.load_image,
            '.gif': self.load_image,
            '.bmp': self.load_image,
            '.webp': self.load_image,
            '.tiff': self.load_image,
            '.tif': self.load_image,
        }
        self._executor = ThreadPoolExecutor(max_workers=4)

    @property
    def supported_formats(self) -> List[str]:
        return list(self._supported_formats.keys())

    def load(self, file_path: str, api_key: str = '', ocr_lang: str = 'chi_sim+eng') -> LoadResult:
        start_time = time.time()
        ext = Path(file_path).suffix.lower()
        if ext not in self._supported_formats:
            return LoadResult(success=False, documents=[], error=f'Unsupported file format: {ext}', file_type=ext, processing_time=time.time()-start_time)
        try:
            loader_func = self._supported_formats[ext]
            if ext == '.pdf':
                docs = loader_func(file_path, api_key=api_key)
            else:
                docs = loader_func(file_path, ocr_lang=ocr_lang)
            char_count = sum(len(d.page_content) for d in docs)
            return LoadResult(success=True, documents=docs, file_type=ext, page_count=len(docs), char_count=char_count, processing_time=time.time()-start_time)
        except Exception as e:
            logger.error(f'Failed to load {file_path}: {e}')
            return LoadResult(success=False, documents=[], error=str(e), file_type=ext, processing_time=time.time()-start_time)

    def load_pdf(self, file_path: str, api_key: str = '') -> List[Document]:
        docs = []
        try:
            import sys
            if 'PyMuPDF' not in sys.modules:
                import pymupdf as fitz
            else:
                import fitz
            doc = fitz.open(file_path)
            total = len(doc)
            for page_num, page in enumerate(doc):
                text = page.get_text()
                if text and text.strip():
                    docs.append(Document(page_content=text.strip(), metadata={'source': file_path, 'page': page_num+1, 'total_pages': total, 'type': 'pdf', 'file_name': Path(file_path).name}))
                else:
                    ocr_text = self._ocr_page(page, api_key)
                    if ocr_text:
                        docs.append(Document(page_content=ocr_text.strip(), metadata={'source': file_path, 'page': page_num+1, 'total_pages': total, 'type': 'pdf_ocr', 'file_name': Path(file_path).name}))
                        logger.info(f'  [OCR p.{page_num+1}/{total}] {len(ocr_text)} chars')
            doc.close()
        except Exception as e:
            logger.error(f'[PDF] {file_path}: {e}')
        return docs

    def _ocr_page(self, page, api_key: str) -> str:
        try:
            from io import BytesIO
            from PIL import Image
            import pytesseract
            pix = page.get_pixmap(dpi=150)
            img_bytes = pix.tobytes('png')
            if api_key:
                try:
                    from ocr import ocr_pdf_page
                    ocr_text = ocr_pdf_page(img_bytes, api_key)
                    if ocr_text and ocr_text.strip():
                        return ocr_text
                except Exception:
                    pass
            img = Image.open(BytesIO(img_bytes))
            return pytesseract.image_to_string(img, lang='chi_sim+eng')
        except Exception as e:
            logger.warning(f'  [OCR failed] {e}')
            return ''

    def load_txt(self, file_path: str, ocr_lang: str = 'chi_sim+eng') -> List[Document]:
        docs = []
        for enc in ['utf-8', 'gbk', 'gb2312', 'latin-1']:
            try:
                with open(file_path, 'r', encoding=enc) as f:
                    content = f.read()
                if content.strip():
                    docs.append(Document(page_content=content.strip(), metadata={'source': file_path, 'page': 1, 'type': 'txt', 'file_name': Path(file_path).name, 'encoding': enc}))
                break
            except (UnicodeDecodeError, LookupError):
                continue
        return docs

    def load_md(self, file_path: str, ocr_lang: str = 'chi_sim+eng') -> List[Document]:
        docs = []
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            if content.strip():
                docs.append(Document(page_content=content.strip(), metadata={'source': file_path, 'page': 1, 'type': 'md', 'file_name': Path(file_path).name}))
        except Exception as e:
            logger.error(f'[MD] {file_path}: {e}')
        return docs

    def load_docx(self, file_path: str, ocr_lang: str = 'chi_sim+eng') -> List[Document]:
        docs = []
        try:
            from docx import Document as DocxDocument
            docx_doc = DocxDocument(file_path)
            paragraphs = [p.text.strip() for p in docx_doc.paragraphs if p.text.strip()]
            full_text = '\n'.join(paragraphs)
            tables_text = []
            for table in docx_doc.tables:
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    if any(cells):
                        tables_text.append(' | '.join(cells))
                tables_text.append('')
            if tables_text:
                full_text += '\n\n[表格内容]\n' + '\n'.join(tables_text)
            if full_text.strip():
                docs.append(Document(page_content=full_text.strip(), metadata={'source': file_path, 'page': 1, 'type': 'docx', 'file_name': Path(file_path).name, 'paragraph_count': len(paragraphs), 'table_count': len(docx_doc.tables)}))
        except ImportError:
            logger.warning('[DOCX] python-docx not installed')
        except Exception as e:
            logger.error(f'[DOCX] {file_path}: {e}')
        return docs

    def load_xlsx(self, file_path: str, ocr_lang: str = 'chi_sim+eng') -> List[Document]:
        docs = []
        try:
            from openpyxl import load_workbook
            wb = load_workbook(file_path, data_only=True)
            all_text = []
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                all_text.append(f'[工作表: {sheet_name}]')
                for row in sheet.iter_rows(values_only=True):
                    cells = [str(cell) if cell is not None else '' for cell in row]
                    if any(cells):
                        all_text.append(' | '.join(cells))
                all_text.append('')
            full_text = '\n'.join(all_text)
            if full_text.strip():
                docs.append(Document(page_content=full_text.strip(), metadata={'source': file_path, 'page': 1, 'type': 'xlsx', 'file_name': Path(file_path).name, 'sheet_count': len(wb.sheetnames), 'sheets': wb.sheetnames}))
        except ImportError:
            logger.warning('[XLSX] openpyxl not installed')
        except Exception as e:
            logger.error(f'[XLSX] {file_path}: {e}')
        return docs

    def load_pptx(self, file_path: str, ocr_lang: str = 'chi_sim+eng') -> List[Document]:
        docs = []
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            all_text = []
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = [f'[幻灯片 {slide_num}]']
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        slide_text.append(shape.text.strip())
                if len(slide_text) > 1:
                    all_text.append('\n'.join(slide_text))
            full_text = '\n\n'.join(all_text)
            if full_text.strip():
                docs.append(Document(page_content=full_text.strip(), metadata={'source': file_path, 'page': 1, 'type': 'pptx', 'file_name': Path(file_path).name, 'slide_count': len(prs.slides)}))
        except ImportError:
            logger.warning('[PPTX] python-pptx not installed')
        except Exception as e:
            logger.error(f'[PPTX] {file_path}: {e}')
        return docs

    def load_csv(self, file_path: str, ocr_lang: str = 'chi_sim+eng') -> List[Document]:
        docs = []
        try:
            import csv, chardet
            with open(file_path, 'rb') as f:
                raw = f.read()
            det = chardet.detect(raw)
            encoding = det.get('encoding') or 'utf-8'
            with open(file_path, 'r', encoding=encoding, errors='replace') as f:
                reader = csv.reader(f)
                rows = list(reader)
            all_text = []
            headers = rows[0] if rows else []
            if headers:
                all_text.append('[CSV Headers] ' + ' | '.join(str(h).strip() for h in headers))
            for row in rows[1:]:
                cells = [str(c).strip() for c in row]
                if any(cells):
                    all_text.append(' | '.join(cells))
            full_text = '\n'.join(all_text)
            if full_text.strip():
                docs.append(Document(page_content=full_text.strip(), metadata={'source': file_path, 'page': 1, 'type': 'csv', 'file_name': Path(file_path).name, 'row_count': max(0, len(rows)-1), 'column_count': len(headers)}))
        except Exception as e:
            logger.error(f'[CSV] {file_path}: {e}')
        return docs

    def load_html(self, file_path: str, ocr_lang: str = 'chi_sim+eng') -> List[Document]:
        docs = []
        try:
            from bs4 import BeautifulSoup
            with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                html = f.read()
            soup = BeautifulSoup(html, 'lxml')
            for tag in soup(['script', 'style', 'nav', 'footer', 'header', 'aside']):
                tag.decompose()
            text = soup.get_text(separator='\n', strip=True)
            text = re.sub(r'\n{3,}', '\n\n', text)
            if text.strip():
                title = ''
                if soup.title and soup.title.string:
                    title = soup.title.string.strip()
                docs.append(Document(page_content=text.strip(), metadata={'source': file_path, 'page': 1, 'type': 'html', 'file_name': Path(file_path).name, 'title': title or Path(file_path).stem}))
        except ImportError:
            logger.warning('[HTML] beautifulsoup4 not installed')
        except Exception as e:
            logger.error(f'[HTML] {file_path}: {e}')
        return docs

    def load_json(self, file_path: str, ocr_lang: str = 'chi_sim+eng') -> List[Document]:
        docs = []
        try:
            import json
            with open(file_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
            text = self._flatten_json(data)
            if text.strip():
                docs.append(Document(page_content=text.strip(), metadata={'source': file_path, 'page': 1, 'type': 'json', 'file_name': Path(file_path).name}))
        except Exception as e:
            logger.error(f'[JSON] {file_path}: {e}')
        return docs

    def _flatten_json(self, obj, prefix: str = '', depth: int = 0) -> str:
        if depth > 10:
            return str(obj)
        parts = []
        if isinstance(obj, dict):
            for k, v in obj.items():
                parts.append(self._flatten_json(v, f'{prefix}{k}: ', depth+1))
        elif isinstance(obj, list):
            for i, v in enumerate(obj):
                parts.append(self._flatten_json(v, f'{prefix}[{i}]: ', depth+1))
        else:
            parts.append(f'{prefix}{obj}')
        return '\n'.join(parts)

    def load_xml(self, file_path: str, ocr_lang: str = 'chi_sim+eng') -> List[Document]:
        docs = []
        try:
            import xml.etree.ElementTree as ET
            tree = ET.parse(file_path)
            root = tree.getroot()
            text = self._flatten_xml(root)
            if text.strip():
                docs.append(Document(page_content=text.strip(), metadata={'source': file_path, 'page': 1, 'type': 'xml', 'file_name': Path(file_path).name, 'root_tag': root.tag}))
        except Exception as e:
            logger.error(f'[XML] {file_path}: {e}')
        return docs

    def _flatten_xml(self, element, prefix: str = '', depth: int = 0) -> str:
        if depth > 10:
            return element.text or ''
        parts = []
        if element.text and element.text.strip():
            parts.append(f'{prefix}{element.tag}: {element.text.strip()}')
        for child in element:
            parts.append(self._flatten_xml(child, prefix + '  ', depth+1))
        return '\n'.join(parts)

    def load_eml(self, file_path: str, ocr_lang: str = 'chi_sim+eng') -> List[Document]:
        docs = []
        try:
            from email import policy
            from email.parser import BytesParser
            with open(file_path, 'rb') as f:
                msg = BytesParser(policy=policy.default).parse(f)
            parts = []
            subject = msg['Subject'] or '(无主题)'
            sender = msg['From'] or '(无发件人)'
            date = msg['Date'] or '(无日期)'
            parts.append(f'[主题: {subject}]')
            parts.append(f'[发件人: {sender}]')
            parts.append(f'[日期: {date}]')
            if msg.is_multipart():
                for part in msg.walk():
                    if part.get_content_type() == 'text/plain':
                        payload = part.get_payload(decode=True)
                        charset = part.get_content_charset() or 'utf-8'
                        try:
                            parts.append(payload.decode(charset, errors='replace'))
                        except Exception:
                            pass
            else:
                payload = msg.get_payload(decode=True)
                charset = msg.get_content_charset() or 'utf-8'
                try:
                    parts.append(payload.decode(charset, errors='replace'))
                except Exception:
                    pass
            full_text = '\n'.join(parts)
            if full_text.strip():
                docs.append(Document(page_content=full_text.strip(), metadata={'source': file_path, 'page': 1, 'type': 'eml', 'file_name': Path(file_path).name, 'subject': subject, 'from': sender}))
        except Exception as e:
            logger.error(f'[EML] {file_path}: {e}')
        return docs

    def load_msg(self, file_path: str, ocr_lang: str = 'chi_sim+eng') -> List[Document]:
        docs = []
        try:
            from extract_msg import Message
            msg_obj = Message(file_path)
            parts = []
            if msg_obj.subject:
                parts.append(f'[主题: {msg_obj.subject}]')
            if msg_obj.sender:
                parts.append(f'[发件人: {msg_obj.sender}]')
            if msg_obj.date:
                parts.append(f'[日期: {msg_obj.date}]')
            body = msg_obj.body
            if body:
                parts.append(body)
            full_text = '\n'.join(parts)
            if full_text.strip():
                docs.append(Document(page_content=full_text.strip(), metadata={'source': file_path, 'page': 1, 'type': 'msg', 'file_name': Path(file_path).name, 'subject': msg_obj.subject or '', 'from': msg_obj.sender or ''}))
        except ImportError:
            logger.warning('[MSG] extract-msg not installed, cannot load .msg files')
        except Exception as e:
            logger.error(f'[MSG] {file_path}: {e}')
        return docs

    def load_image(self, file_path: str, ocr_lang: str = 'chi_sim+eng') -> List[Document]:
        docs = []
        try:
            from PIL import Image
            import pytesseract
            image = Image.open(file_path)
            text = pytesseract.image_to_string(image, lang=ocr_lang)
            if text and text.strip():
                docs.append(Document(page_content=text.strip(), metadata={'source': file_path, 'page': 1, 'type': 'image_ocr', 'file_name': Path(file_path).name, 'image_size': image.size, 'image_mode': image.mode}))
        except Exception as e:
            logger.error(f'[OCR] {file_path}: {e}')
        return docs


def load_document(file_path: str, ocr_lang: str = 'chi_sim+eng', api_key: str = '') -> List[Document]:
    loader = DocumentLoader()
    result = loader.load(file_path, api_key=api_key, ocr_lang=ocr_lang)
    return result.documents

def load_folder(folder_path: str, ocr_lang: str = 'chi_sim+eng', api_key: str = '') -> Tuple[List[Document], List[Dict[str, Any]]]:
    loader = DocumentLoader()
    all_docs = []
    results = []
    for root, _, files in os.walk(folder_path):
        for fname in files:
            ext = Path(fname).suffix.lower()
            if ext in loader.supported_formats:
                fpath = os.path.join(root, fname)
                result = loader.load(fpath, api_key=api_key, ocr_lang=ocr_lang)
                all_docs.extend(result.documents)
                results.append({'file': fname, 'path': fpath, **result.to_dict()})
    return all_docs, results

load_pdf = lambda fp, api_key='': DocumentLoader().load_pdf(fp, api_key)
load_txt = lambda fp, ocr_lang='chi_sim+eng': DocumentLoader().load_txt(fp, ocr_lang)
load_md = lambda fp, ocr_lang='chi_sim+eng': DocumentLoader().load_md(fp, ocr_lang)
load_image = lambda fp, ocr_lang='chi_sim+eng': DocumentLoader().load_image(fp, ocr_lang)
